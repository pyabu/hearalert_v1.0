#!/usr/bin/env python3
"""
Generate HearAlert Zeroth Review PowerPoint Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RgbColor(31, 73, 125)
LIGHT_BLUE = RgbColor(79, 129, 189)
WHITE = RgbColor(255, 255, 255)
BLACK = RgbColor(0, 0, 0)
GRAY = RgbColor(100, 100, 100)
GREEN = RgbColor(46, 125, 50)

def add_title_slide(prs, title, subtitle, author, date):
    """Add title slide with blue background"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.2), Inches(2.333), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = RgbColor(150, 150, 150)
    line.line.fill.background()
    
    # Author
    auth_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.333), Inches(1.5))
    tf = auth_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = author
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = date
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    """Add content slide with bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Title underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.333), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)
        p.font.color.rgb = BLACK
        p.space_before = Pt(10)
    
    # Bottom accent line
    bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.4), prs.slide_width, Pt(8))
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = LIGHT_BLUE
    bottom_line.line.fill.background()

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add two-column content slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Title underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.333), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_before = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.1), Inches(6), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_before = Pt(8)
    
    # Bottom accent line
    bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.4), prs.slide_width, Pt(8))
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = LIGHT_BLUE
    bottom_line.line.fill.background()

# ============ CREATE SLIDES ============

# Slide 1: Title
add_title_slide(
    prs,
    "HearAlert: AI-Powered Real-Time Sound Recognition\nand Multi-Modal Alert System",
    "For Deaf and Hard-of-Hearing Individuals",
    "S. Abusaleem (24MSCSCPY0037)\nM.Sc. Computer Science\nDepartment of Computer Science\nPondicherry University",
    "22 January 2026"
)

# Slide 2: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "430+ million people worldwide have disabling hearing loss (WHO)",
    "Safety Risks: Unable to hear fire alarms, vehicle horns, sirens",
    "Communication Barriers: Missing doorbells, phone calls, interactions",
    "Childcare Challenges: Cannot detect baby crying or infant distress",
    "Limited Solutions: Traditional hearing aids don't cover all environments",
    "Solution: AI-powered mobile app with multi-sensory alerts"
])

# Slide 3: Solution
add_two_column_slide(prs, "HearAlert Solution",
    "🎤 Sound Detection", [
        "Real-time audio capture at 16kHz",
        "TensorFlow Lite + YAMNet model",
        "500+ sound categories",
        "Less than 50ms latency"
    ],
    "📱 Multi-Modal Alerts", [
        "Haptic: Category-specific vibrations",
        "Visual: Flashlight strobe patterns",
        "Audio: Text-to-speech announcements",
        "Screen: Full-screen visual alerts"
    ]
)

# Slide 4: Technology Stack
add_content_slide(prs, "Technology Stack", [
    "Framework: Flutter 3.24+ (Cross-platform)",
    "Language: Dart 3.5.3+",
    "ML Runtime: TensorFlow Lite",
    "Audio Model: YAMNet (521 classes, AudioSet trained)",
    "Audio Capture: mic_stream package (16kHz)",
    "State Management: Provider Pattern",
    "Platforms: Android 5.0+, iOS 12.0+"
])

# Slide 5: Key Features
add_two_column_slide(prs, "Key Features",
    "🔊 Priority Sound Detection", [
        "Fire alarm / Smoke detector",
        "Vehicle horn / Sirens",
        "Door knock / Doorbell",
        "Baby cry / Infant distress",
        "Glass breaking / Danger sounds"
    ],
    "📳 Smart Vibration Patterns", [
        "Fire Alarm: SOS pattern (... --- ...)",
        "Vehicle Horn: Long warning pulses",
        "Door Knock: Double tap pattern",
        "Baby Cry: Gentle pulse",
        "Glass Break: Sharp urgent jitter"
    ]
)

# Slide 6: Advanced Features
add_two_column_slide(prs, "Advanced Features",
    "👶 Baby Cry Classifier", [
        "Hungry: Try feeding",
        "Burping: Pat back gently",
        "Pain: Tummy massage",
        "Discomfort: Check diaper",
        "Tired: Calm environment"
    ],
    "🏠 Scenario Profiles & SOS", [
        "Home: Doorbell, baby, appliances",
        "Street: Horns, sirens, traffic",
        "School: Bells, announcements",
        "Emergency SOS: One-tap trigger"
    ]
)

# Slide 7: Plan of Work
add_content_slide(prs, "Plan of Work", [
    "Phase 1: Research & Planning ✓ Complete",
    "Phase 2: Core Development (Audio, ML) ○ Planned",
    "Phase 3: Alert System (Vibration, Flash, TTS) ○ Planned",
    "Phase 4: UI/UX Development ○ Planned",
    "Phase 5: Advanced Features ○ Planned",
    "Phase 6: Testing & Deployment ○ Planned"
])

# Slide 8: System Requirements
add_two_column_slide(prs, "System Requirements",
    "📱 Target Devices", [
        "Android: 5.0+ (API 21)",
        "iOS: 12.0+",
        "RAM: 2 GB minimum",
        "Storage: 100 MB"
    ],
    "🔧 Hardware Features", [
        "Microphone: Required",
        "Vibration Motor: Required",
        "Camera Flash: Recommended",
        "Internet: Not Required"
    ]
)

# Slide 9: Thank You
add_title_slide(
    prs,
    "Thank You",
    "Questions?",
    "S. Abusaleem (24MSCSCPY0037)\nM.Sc. Computer Science\nDepartment of Computer Science\nPondicherry University",
    ""
)

# Save
prs.save('/home/red-dragon/Desktop/my/Adaptive-Puzzles/mobile_app/HEARALERT_PRESENTATION.pptx')
print("✅ Presentation saved: HEARALERT_PRESENTATION.pptx")
